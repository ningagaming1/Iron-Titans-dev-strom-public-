"""
build_deck.py - builds docs/SmartHome.pptx (~50 slides).

    .venv/bin/python docs/build_deck.py

Plain python-pptx. Look matches the app: paper+ink, terracotta accent,
forest-green headings.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "SmartHome.pptx")

# palette, from web/style.css
PAPER      = RGBColor(0xFD, 0xF9, 0xEF)
PAPER_2    = RGBColor(0xF1, 0xEA, 0xD7)
INK        = RGBColor(0x2B, 0x27, 0x1F)
INK_SOFT   = RGBColor(0x6F, 0x68, 0x57)
FOREST     = RGBColor(0x31, 0x3D, 0x2D)
FOREST_2   = RGBColor(0x3D, 0x4B, 0x38)
ACCENT     = RGBColor(0xB2, 0x5C, 0x3F)
ACCENT_D   = RGBColor(0x94, 0x47, 0x2E)
GOOD       = RGBColor(0x4F, 0x6B, 0x43)
DARK       = RGBColor(0x23, 0x21, 0x1B)
DARK_2     = RGBColor(0x2E, 0x2A, 0x22)
CREAM      = RGBColor(0xF1, 0xEA, 0xD7)
LINE       = RGBColor(0xDD, 0xD0, 0xB4)

HEAD_FONT = "Georgia"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
MARGIN = Inches(0.92)
CONTENT_W = EMU_W - 2 * MARGIN

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_num = 0


# --- helpers ---
def slide(dark=False):
    global _num
    _num += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, EMU_W, EMU_H, DARK if dark else PAPER, line=None)
    return s


def rect(s, x, y, w, h, fill, line=None, line_w=1.0, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '38100', 'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '2B271F'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '22000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp


def rrect(s, x, y, w, h, fill, line=None, line_w=1.0, shadow=False):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        sp.adjustments[0] = 0.06
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = el.makeelement(qn('a:outerShdw'),
                            {'blurRad': '80000', 'dist': '30000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = el.makeelement(qn('a:srgbClr'), {'val': '2B271F'})
        alpha = el.makeelement(qn('a:alpha'), {'val': '20000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    return sp


def text(s, x, y, w, h, runs, size=18, color=INK, bold=False, italic=False,
         font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.12, space_after=6):
    """runs: a string, or list of (string, dict-overrides)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        if isinstance(item, tuple):
            txt, ov = item
        else:
            txt, ov = item, {}
        r = p.add_run()
        r.text = txt
        f = r.font
        f.name = ov.get("font", font)
        f.size = Pt(ov.get("size", size))
        f.bold = ov.get("bold", bold)
        f.italic = ov.get("italic", italic)
        f.color.rgb = ov.get("color", color)
        try:
            f.language_id = MSO_LANGUAGE_ID.ENGLISH_US
        except Exception:
            pass
    return tb


def bullets(s, items, x=None, y=Inches(1.95), w=None, size=17, gap=10,
            color=INK, dark=False):
    x = x if x is not None else MARGIN
    w = w if w is not None else CONTENT_W
    tb = s.shapes.add_textbox(x, y, w, EMU_H - y - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap if lvl == 0 else 4)
        p.line_spacing = 1.13
        r = p.add_run()
        r.text = ("- " if lvl else "") + it if lvl else it
        r.font.name = BODY_FONT
        r.font.size = Pt(size if lvl == 0 else size - 3)
        r.font.color.rgb = (CREAM if dark else color) if lvl == 0 else \
            (RGBColor(0xB9, 0xB0, 0x9C) if dark else INK_SOFT)
        r.font.bold = (lvl == 0 and False)
        # the bullet marker
        pPr = p._p.get_or_add_pPr()
        for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum'):
            for e in pPr.findall(qn(tag)):
                pPr.remove(e)
        if lvl == 0:
            buClr = pPr.makeelement(qn('a:buClr'), {})
            c = pPr.makeelement(qn('a:srgbClr'), {'val': 'B25C3F'})
            buClr.append(c); pPr.append(buClr)
            buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
            buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
            pPr.append(buFont); pPr.append(buChar)
        else:
            buNone = pPr.makeelement(qn('a:buNone'), {})
            pPr.append(buNone)
    return tb


def title(s, t, kicker=None, dark=False):
    rect(s, MARGIN, Inches(0.66), Inches(0.09), Inches(0.62), ACCENT)
    if kicker:
        text(s, MARGIN + Inches(0.28), Inches(0.5), CONTENT_W, Inches(0.3),
             kicker.upper(), size=11.5, color=ACCENT_D, bold=True, font=MONO_FONT)
        ty = Inches(0.78)
    else:
        ty = Inches(0.62)
    text(s, MARGIN + Inches(0.28), ty, CONTENT_W, Inches(0.9), t,
         size=30, color=(CREAM if dark else FOREST), bold=True, font=HEAD_FONT)
    footer(s, dark)


TOTAL = 59  # fixed up on the second pass


def footer(s, dark=False):
    text(s, MARGIN, EMU_H - Inches(0.5), Inches(4), Inches(0.3),
         "SmartHome  -  Iron Titans", size=9.5,
         color=(RGBColor(0x8A, 0x82, 0x70) if dark else INK_SOFT), font=MONO_FONT)
    text(s, EMU_W - MARGIN - Inches(1.4), EMU_H - Inches(0.5), Inches(1.4), Inches(0.3),
         f"{_num:02d} / {TOTAL}", size=9.5, align=PP_ALIGN.RIGHT,
         color=(RGBColor(0x8A, 0x82, 0x70) if dark else INK_SOFT), font=MONO_FONT)


def box(s, x, y, w, h, label, fill=CREAM, fg=INK, size=12.5, line=LINE, bold=False):
    sp = rrect(s, x, y, w, h, fill, line=line, line_w=1.2)
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.name = BODY_FONT; r.font.size = Pt(size); r.font.color.rgb = fg; r.font.bold = bold
    return sp


def connect(s, x1, y1, x2, y2, color=ACCENT_D, w=1.6, dash=False):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    le = cn.line._get_or_add_ln()
    # a:ln children must stay in schema order: prstDash before tailEnd
    if dash:
        d = le.makeelement(qn('a:prstDash'), {'val': 'dash'})
        le.append(d)
    tail = le.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    le.append(tail)
    return cn


def code_slide(kicker, ttl, code, note=None):
    s = slide(dark=True)
    title(s, ttl, kicker, dark=True)
    panel = rrect(s, MARGIN, Inches(1.85), CONTENT_W,
                  EMU_H - Inches(1.85) - Inches(0.95) - (Inches(0.5) if note else 0),
                  DARK_2, line=RGBColor(0x47, 0x40, 0x33), line_w=1.0)
    tf = panel.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.32); tf.margin_right = Inches(0.32)
    tf.margin_top = Inches(0.24); tf.margin_bottom = Inches(0.24)
    for i, ln in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.16
        r = p.add_run(); r.text = ln if ln else " "
        r.font.name = MONO_FONT; r.font.size = Pt(13)
        low = ln.strip()
        if low.startswith("#") or low.startswith("//"):
            r.font.color.rgb = RGBColor(0x9A, 0x8F, 0x76)
        elif any(k in ln for k in ("def ", "class ", "return", "import ", "for ", "if ", "with ")):
            r.font.color.rgb = RGBColor(0xE7, 0xC9, 0x8C)
        else:
            r.font.color.rgb = RGBColor(0xEC, 0xE4, 0xD3)
    if note:
        text(s, MARGIN, EMU_H - Inches(1.32), CONTENT_W, Inches(0.5), note,
             size=12.5, italic=True, color=RGBColor(0xB9, 0xB0, 0x9C), font=BODY_FONT)
    return s


def section(kicker, ttl, sub=None):
    s = slide(dark=True)
    rect(s, 0, Inches(3.15), EMU_W, Inches(0.02), RGBColor(0x47, 0x40, 0x33))
    rect(s, MARGIN, Inches(2.55), Inches(0.6), Inches(0.1), ACCENT)
    text(s, MARGIN, Inches(2.75), CONTENT_W, Inches(0.4), kicker.upper(),
         size=13, color=ACCENT, bold=True, font=MONO_FONT)
    text(s, MARGIN, Inches(3.25), CONTENT_W, Inches(1.2), ttl,
         size=40, color=CREAM, bold=True, font=HEAD_FONT)
    if sub:
        text(s, MARGIN, Inches(4.35), CONTENT_W, Inches(0.8), sub,
             size=15, color=RGBColor(0xB9, 0xB0, 0x9C), italic=True, font=BODY_FONT)
    footer(s, dark=True)
    return s


def content(kicker, ttl, items, size=17, gap=10):
    s = slide()
    title(s, ttl, kicker)
    bullets(s, items, size=size, gap=gap)
    return s


def two_col(kicker, ttl, left_ttl, left, right_ttl, right):
    s = slide()
    title(s, ttl, kicker)
    colw = (CONTENT_W - Inches(0.6)) / 2
    for i, (ct, items) in enumerate([(left_ttl, left), (right_ttl, right)]):
        cx = MARGIN + i * (colw + Inches(0.6))
        rect(s, cx, Inches(2.0), Inches(0.5), Inches(0.07), ACCENT)
        text(s, cx, Inches(2.15), colw, Inches(0.4), ct, size=15, bold=True,
             color=FOREST_2, font=HEAD_FONT)
        bullets(s, items, x=cx, y=Inches(2.7), w=colw, size=14, gap=8)
    return s


def table_slide(kicker, ttl, headers, rows, col_ratios=None):
    s = slide()
    title(s, ttl, kicker)
    n = len(headers)
    ratios = col_ratios or [1] * n
    tot = sum(ratios)
    x0, y0 = MARGIN, Inches(2.0)
    tw = CONTENT_W
    rh = min(Inches(0.62), (EMU_H - y0 - Inches(0.9)) / (len(rows) + 1))
    # header
    cx = x0
    for j, hdr in enumerate(headers):
        cw = Emu(int(tw * ratios[j] / tot))
        rect(s, cx, y0, cw, rh, FOREST)
        text(s, cx + Inches(0.12), y0, cw - Inches(0.2), rh, hdr, size=12.5,
             color=CREAM, bold=True, anchor=MSO_ANCHOR.MIDDLE, font=BODY_FONT)
        cx += cw
    for i, row in enumerate(rows):
        ry = y0 + rh * (i + 1)
        cx = x0
        band = PAPER if i % 2 == 0 else PAPER_2
        for j, cell in enumerate(row):
            cw = Emu(int(tw * ratios[j] / tot))
            rect(s, cx, ry, cw, rh, band, line=LINE, line_w=0.75)
            text(s, cx + Inches(0.12), ry, cw - Inches(0.2), rh, cell, size=11.5,
                 color=INK, anchor=MSO_ANCHOR.MIDDLE, font=(MONO_FONT if j == 0 else BODY_FONT))
            cx += cw
    return s


def stat_slide(kicker, ttl, stats):
    s = slide()
    title(s, ttl, kicker)
    n = len(stats)
    gap = Inches(0.5)
    cw = (CONTENT_W - gap * (n - 1)) / n
    for i, (big, small) in enumerate(stats):
        cx = MARGIN + i * (cw + gap)
        rrect(s, cx, Inches(2.6), cw, Inches(2.3), CREAM, line=LINE, line_w=1.0, shadow=True)
        text(s, cx, Inches(2.95), cw, Inches(1.1), big, size=44, bold=True,
             color=ACCENT, align=PP_ALIGN.CENTER, font=HEAD_FONT)
        text(s, cx + Inches(0.2), Inches(4.05), cw - Inches(0.4), Inches(0.8), small,
             size=12.5, color=INK_SOFT, align=PP_ALIGN.CENTER, font=BODY_FONT)
    return s


# --- the deck ---

# ---- 1. cover ----
s = slide(dark=True)
rect(s, 0, 0, Inches(0.18), EMU_H, ACCENT)
text(s, MARGIN, Inches(1.5), CONTENT_W, Inches(0.5), "HACKATHON BUILD  -  IRON TITANS",
     size=13, color=ACCENT, bold=True, font=MONO_FONT)
text(s, MARGIN, Inches(2.2), CONTENT_W, Inches(1.6),
     [("SmartHome", {})], size=68, color=CREAM, bold=True, font=HEAD_FONT)
text(s, MARGIN, Inches(3.6), CONTENT_W, Inches(0.9),
     "The house is on the guest list. Control the lights, the fan and the front "
     "door by click, by type, or by voice - offline and free.",
     size=17, color=RGBColor(0xC7, 0xBE, 0xAA), italic=True, font=BODY_FONT)
text(s, MARGIN, Inches(5.5), CONTENT_W, Inches(0.9),
     [("Team  ", {"color": RGBColor(0x8A, 0x82, 0x70), "font": MONO_FONT, "size": 12}),
      ("Pranjal  -  Samarpreet  -  Ekamjot  -  Amrit", {"color": CREAM, "size": 14, "bold": True})])
text(s, MARGIN, Inches(6.05), CONTENT_W, Inches(0.5),
     "Python standard library  -  Vosk  -  Piper  -  vanilla JS", size=12,
     color=RGBColor(0x8A, 0x82, 0x70), font=MONO_FONT)
footer(s, dark=True)

# ---- 2. elevator pitch --------------------------------------
s = slide(dark=True)
text(s, MARGIN, Inches(1.9), CONTENT_W, Inches(1.2),
     [("One page. Three ways in.", {"color": CREAM, "size": 36, "bold": True, "font": HEAD_FONT})])
text(s, MARGIN, Inches(3.15), CONTENT_W - Inches(1.2), Inches(3),
     [("Tap a button, type a command, or just speak.  ",
       {"color": CREAM, "size": 21}),
      ("The recognizer, the reply voice and the logic all run on your own "
       "machine - no cloud, no API keys, no bill.",
       {"color": RGBColor(0xC7, 0xBE, 0xAA), "size": 21})],
     line_spacing=1.3)
footer(s, dark=True)

# ---- 3. agenda ---------------------------------------------
content("start here", "What we'll walk through", [
    "The problem - home control is scattered, voice assistants are rented",
    "The build - one stdlib server, JSON for state, a browser front-end",
    "Accounts - request, admin approval, scrambled passwords, lockout",
    "The house - devices, activity log, and the intent parser",
    "Voice - offline speech in (Vosk) and out (Piper), tuned for accuracy",
    "The front-end - dashboard, dark mode, the listening pop-up",
    "Quality - the self-test suite, what we fixed, what's left",
    "Roadmap - real hardware, a wake word, multi-home",
], size=16, gap=9)

# ---- 4. section: problem ----------------------------------
section("01", "The problem", "Why build another home dashboard?")

# ---- 5 ----
content("problem", "Home control is fragmented", [
    "Every bulb, plug and lock ships its own app",
    "Three rooms can mean four apps and four logins",
    "No single 'what is on right now' view",
    "Guests can't help - sharing means handing over an account",
    "The simple ask - 'turn the light off' - takes six taps",
])

# ---- 6 ----
content("problem", "Voice assistants are rented, not owned", [
    "The big assistants stream your microphone to a company server",
    "That needs an internet connection just to hear 'lights off'",
    "Free tiers have quotas; the good tiers cost money",
    "Your commands - and your home's state - live on someone else's box",
    "Offline hobby setups exist but are fiddly to stand up",
])

# ---- 7 ----
two_col("solution", "What we built", "The app", [
    "One small web server (Python stdlib only)",
    "A dashboard for light / fan / door",
    "Admin-gated sign-up",
    "Shared house state, live on every screen",
], "The voice layer", [
    "Speech - text with Vosk, on-device",
    "Text - speech with Piper, natural voice",
    "Tuned to the command vocabulary",
    "Falls back to the browser engine if models are absent",
])

# ---- 8 ----
content("principles", "Four rules we held to", [
    "Free - every dependency is open source, no keys, no quota",
    "Offline - after a one-time model download, nothing phones home",
    "Small - the core runs on the standard library; no framework, no database",
    "Safe by default - new accounts wait for an admin; passwords are scrambled",
], size=17, gap=14)

# ---- 9 ----
content("demo", "What you can actually do", [
    "\"turn on the light\"  /  \"switch off the fan\"",
    "\"lock the door\"  /  \"open the door\"",
    "\"turn on everything\"  /  \"turn everything off\"",
    "Type the same phrases, or use the device buttons and quick controls",
    "Watch the activity log fill in - and the room go dark when the light is off",
])

# ---- 10. section: architecture ---------------------------
section("02", "Architecture", "One server, some JSON files, a browser")

# ---- 11. architecture diagram ----
s = slide()
title(s, "How the pieces fit", "architecture")
box(s, MARGIN, Inches(2.15), Inches(3.2), Inches(1.15), "Browser\nweb/ - HTML, CSS, JS", CREAM, INK, 12.5, bold=True)
box(s, MARGIN, Inches(4.0), Inches(3.2), Inches(1.15), "Microphone\nMediaRecorder", PAPER_2, INK, 12.5)
box(s, Inches(5.4), Inches(2.15), Inches(3.0), Inches(3.0),
    "app.py\n\nstdlib http.server\nserves web/ + JSON API", FOREST, CREAM, 13, bold=True)
box(s, Inches(9.5), Inches(1.55), Inches(2.9), Inches(0.95), "login2.py\nsignup.py", CREAM, INK, 12)
box(s, Inches(9.5), Inches(2.75), Inches(2.9), Inches(0.95), "devices.py\nintent.py", CREAM, INK, 12)
box(s, Inches(9.5), Inches(3.95), Inches(2.9), Inches(0.95), "voice.py\nVosk + Piper", CREAM, INK, 12)
box(s, Inches(9.5), Inches(5.15), Inches(2.9), Inches(0.8), "data/*.json\nusers, house", PAPER_2, INK, 11.5)
connect(s, Inches(4.12), Inches(2.7), Inches(5.4), Inches(3.0))
connect(s, Inches(4.12), Inches(4.55), Inches(5.4), Inches(3.9))
connect(s, Inches(8.4), Inches(2.9), Inches(9.5), Inches(2.0))
connect(s, Inches(8.4), Inches(3.3), Inches(9.5), Inches(3.2))
connect(s, Inches(8.4), Inches(3.9), Inches(9.5), Inches(4.4))
connect(s, Inches(10.95), Inches(4.9), Inches(10.95), Inches(5.15), color=INK_SOFT)
text(s, MARGIN, Inches(5.7), Inches(7), Inches(0.8),
     "No framework. No ORM. No build step. Open the folder, run one file.",
     size=13, italic=True, color=INK_SOFT)

# ---- 12. stack ----
table_slide("stack", "The whole toolbox", ["Layer", "Choice", "Why"], [
    ["Server", "http.server (stdlib)", "Zero install, easy to read, enough for a demo"],
    ["State", "JSON files", "Inspectable, diff-able, no DB to run"],
    ["Front-end", "Vanilla HTML/CSS/JS", "No build, no npm, loads instantly"],
    ["STT", "Vosk (small EN model)", "Offline, fast, ~40 MB, grammar support"],
    ["TTS", "Piper (neural)", "Offline, natural voice, tiny + quick"],
    ["Audio glue", "ffmpeg", "Decodes the mic blob, shifts TTS pitch"],
], col_ratios=[1.1, 1.6, 3.2])

# ---- 13. repo layout ----
table_slide("layout", "Every file, one job", ["File", "Responsibility"], [
    ["main.py", "Run everything: checks, seed, self-test, start server"],
    ["app.py", "HTTP server - static files + the JSON API"],
    ["login2.py", "Accounts: request, approve/reject, login, lockout"],
    ["signup.py", "Password scrambling (hash rounds + Fibonacci check)"],
    ["devices.py", "House state (light/fan/door) + activity log"],
    ["intent.py", "Plain text -> a structured command"],
    ["voice.py", "Offline speech in (Vosk) and out (Piper)"],
    ["voice_setup.py", "One-time model download"],
    ["seed.py", "Reset databases, create the starter admin"],
    ["web/", "index.html, dashboard.html, auth.js, dashboard.js, style.css"],
], col_ratios=[1.2, 4])

# ---- 14. data model ----
two_col("data", "State is just three JSON files", "data/users/", [
    "users.json - approved accounts",
    "pending.json - waiting for approval",
    "settings.json - dev-mode toggle (git-ignored)",
    "Each user: rounds, fib_check, password hash, lock state",
], "data/", [
    "devices.json - the one shared house",
    "light: bool, fan: bool, door_locked: bool",
    "activity: last 12 actions with timestamps",
    "models/ - downloaded Vosk + Piper (git-ignored)",
])

# ---- 15. request lifecycle ----
s = slide()
title(s, "A request, start to finish", "lifecycle")
steps = ["Browser\nfetch()", "app.py\ndo_GET / do_POST", "read JSON\nbody", "module\nlogin2 / devices / voice", "write\ndata/*.json", "JSON\nresponse"]
bw = Inches(1.9); gap = Inches(0.12); y = Inches(3.0)
x = MARGIN
for i, st in enumerate(steps):
    box(s, x, y, bw, Inches(1.3), st, CREAM if i not in (1, 3) else FOREST,
        INK if i not in (1, 3) else CREAM, 11, bold=(i in (1, 3)))
    if i < len(steps) - 1:
        connect(s, x + bw, y + Inches(0.65), x + bw + gap, y + Inches(0.65))
    x += bw + gap
text(s, MARGIN, Inches(4.7), CONTENT_W, Inches(1),
     "Threaded server, so two laptops can poke the same house at once. "
     "Every write is a full-file JSON dump - simple, and fine at this scale.",
     size=13, italic=True, color=INK_SOFT)

# ---- 16. API ----
table_slide("api", "The API surface", ["Method + path", "Does"], [
    ["POST /api/signup", "Create a pending account request"],
    ["POST /api/login", "Check credentials, return a safe user record"],
    ["POST /api/pending | approve | reject", "Admin manages the waiting list"],
    ["POST /api/devmode", "Toggle auto-approval"],
    ["GET  /api/devices", "The whole house state"],
    ["POST /api/devices/set", "Flip one device (or all)"],
    ["POST /api/command", "Typed text -> intent -> house"],
    ["GET  /api/voice/status", "Is the offline engine ready?"],
    ["POST /api/voice", "Raw mic audio -> text -> house -> spoken reply"],
    ["POST /api/voice/tts", "Text -> a WAV of the reply"],
], col_ratios=[2.1, 3])

# ---- 17. design rationale ----
two_col("rationale", "Why no framework, why no database", "What we skipped", [
    "Flask / FastAPI - not needed for ~12 routes",
    "React / a bundler - the pages are small and static",
    "SQLite / Postgres - state is tiny and rarely written",
    "Docker - it's 'python main.py'",
], "What we got", [
    "A new teammate reads the whole codebase in an hour",
    "git diff shows exactly what changed in the 'database'",
    "Nothing to install for the core app",
    "The self-test can copy data/ to a temp dir and go wild",
])

# ---- 18. section: auth ----
section("03", "Accounts", "Request -> an admin waves you in -> you're in")

# ---- 19. auth flow ----
s = slide()
title(s, "The approval flow", "accounts")
box(s, MARGIN, Inches(2.4), Inches(2.7), Inches(1.2), "Sign up\nrequest_account()", CREAM, INK, 12)
box(s, Inches(4.5), Inches(2.4), Inches(2.7), Inches(1.2), "pending.json\nwaiting", PAPER_2, INK, 12)
box(s, Inches(8.3), Inches(2.4), Inches(2.7), Inches(1.2), "users.json\napproved = admin", FOREST, CREAM, 12, bold=True)
box(s, Inches(4.5), Inches(4.5), Inches(2.7), Inches(1.1), "reject()\ndropped", PAPER_2, INK_SOFT, 11.5)
box(s, Inches(8.3), Inches(4.5), Inches(2.7), Inches(1.1), "login()\nsession in localStorage", CREAM, INK, 11.5)
connect(s, Inches(3.4), Inches(3.0), Inches(4.5), Inches(3.0))
connect(s, Inches(7.2), Inches(3.0), Inches(8.3), Inches(3.0), color=GOOD)
connect(s, Inches(5.8), Inches(3.6), Inches(5.8), Inches(4.5), color=INK_SOFT, dash=True)
connect(s, Inches(9.65), Inches(3.6), Inches(9.65), Inches(4.5), color=GOOD)
text(s, MARGIN, Inches(5.7), CONTENT_W, Inches(0.9),
     "Developer mode short-circuits the wait: sign-ups land straight in users.json. "
     "Handy while building, one toggle on the Invites tab.", size=13, italic=True, color=INK_SOFT)

# ---- 20. two databases ----
content("accounts", "Two files, one direction of travel", [
    "pending.json holds people who signed up but can't log in yet",
    "users.json holds approved people - and in this build, approved means admin",
    "approve() pops the record from one file and writes it to the other",
    "login() only ever reads users.json",
    "No 'roles' table, no join - just which file your name is in",
])

# ---- 21. password scrambling ----
code_slide("signup.py", "Passwords are never stored raw", '''def password_funct(user_password, rounds=None):
    if rounds is None:
        rounds = random.randint(4, 100)      # per-user, saved like a salt

    password = user_password
    for _ in range(rounds):                  # SHA-256, rounds times
        password = hashlib.sha256(password.encode()).hexdigest()

    a, b = 1, 1                               # walk Fibonacci `rounds` steps
    for _ in range(rounds - 1):
        a, b = b, a + b
    fib_check = hashlib.sha256(str(a).encode()).hexdigest()

    return (rounds, fib_check, password)''',
    note="We store (rounds, fib_check, hash). Login re-runs this with the saved "
         "rounds and compares - match on both fingerprints or it's rejected.")

# ---- 22. rounds + fib ----
content("accounts", "Why rounds, why Fibonacci", [
    "rounds is random per user and saved - so the same password hashes "
    "differently for two people (a salt, essentially)",
    "It also makes each check do 4-100 SHA-256 passes - a small, deliberate cost",
    "The Fibonacci hash is a second fingerprint tied to the same rounds number",
    "Both must match on login - a tiny bit of defence in depth",
    "Honest note: this is a learning exercise, not bcrypt / argon2",
])

# ---- 23. login rules ----
two_col("accounts", "Login is a small rulebook", "Rejections", [
    "Unknown username -> generic 'wrong username or password'",
    "Still in pending.json -> 'waiting for approval'",
    "5 wrong tries -> account locks",
], "Ways back in", [
    "Correct password -> counter resets, lock clears, last_login stamped",
    "Locked? paste the stored hash from users.json as the password",
    "Or an admin calls unlock()",
])

# ---- 24. dev mode ----
content("accounts", "Developer mode", [
    "OFF (default) - sign-ups wait for an admin on the Invites tab",
    "ON - sign-ups are auto-approved and can log in immediately",
    "Source of truth: SMARTHOME_DEV_MODE env var, else settings.json, else default",
    "Flippable at runtime: the Invites tab, or POST /api/devmode",
    "The login page shows a badge and relaxes the password minimum when it's on",
])

# ---- 25. admin ----
content("accounts", "The Invites tab", [
    "Lists everyone in pending.json with Accept / Decline buttons",
    "A badge shows the count; it polls every 8 seconds",
    "Accept -> approve(); the person can log in on their next try",
    "Decline -> reject(); the request is gone",
    "The developer-mode switch lives here too",
])

# ---- 26. section: devices ----
section("04", "The house & the parser", "From 'turn off the fan' to a flipped switch")

# ---- 27. house model ----
two_col("devices", "The house is a dict", "Devices", [
    "light - on / off",
    "fan - on / off",
    "door_locked - True means locked",
], "Activity", [
    "A list of the last 12 actions",
    "Each: who did what, an icon key, a HH:MM stamp",
    "New entries go on top; the tail is trimmed",
    "The dashboard renders it as a feed",
])

# ---- 28. devices.py ----
code_slide("devices.py", "Every change goes through one path", '''def set_device(device, value, who="someone"):
    house = _load()                      # read data/devices.json
    value = bool(value)

    if device == "light":
        house["light"] = value
        _log(house, f"{who} turned the light {'on' if value else 'off'}", "light")
    elif device == "fan":
        ...
    elif device == "door":
        house["door_locked"] = value
        _log(house, f"{who} {'locked' if value else 'unlocked'} the door", ...)

    _save(house)                         # write it back, whole file
    return True, house''',
    note="all_devices() and apply() are thin wrappers over this. One writer, "
         "one log call, one save.")

# ---- 29. intent.py ----
code_slide("intent.py", "The parser has no ML in it", '''DEVICE_WORDS  = {"light": "light", "lamp": "light", "bulb": "light",
                 "fan": "fan", "cooler": "fan",
                 "door": "door", "gate": "door", "lock": "door"}
TURN_ON_WORDS  = {"on", "start", "begin", "enable", "activate"}
TURN_OFF_WORDS = {"off", "stop", "disable", "shut", "kill"}

def parse(text):
    words = re.sub(r"[^a-z\\s]", " ", text.lower()).split()
    targets = [DEVICE_WORDS[w] for w in words if w in DEVICE_WORDS]
    if word_set & EVERYTHING_WORDS:
        targets = ["light", "fan"]
    # door has its own lock/unlock words; light+fan look for on vs off
    ...''',
    note="Keyword sets, a few rules, and a friendly 'say' string for every "
         "outcome - including the failures.")

# ---- 30. intent examples ----
table_slide("intent", "parse() in practice", ["You say", "It returns"], [
    ["turn on the light", "action=on, targets=[light]"],
    ["switch off the fan", "action=off, targets=[fan]"],
    ["open the door", "action=unlock, targets=[door]"],
    ["lock the door", "action=lock, targets=[door]"],
    ["turn on everything", "action=on, targets=[light, fan]"],
    ["banana", "ok=False, say='Which device?'"],
], col_ratios=[1.6, 2.6])

# ---- 31. apply ----
content("devices", "intent -> action", [
    "devices.apply(parsed, who) takes the dict from intent.parse()",
    "action 'lock' / 'unlock' -> set the door",
    "action 'on' / 'off' -> loop the targets, flip light and/or fan",
    "Not understood -> change nothing, return the friendly 'say' line",
    "Always returns (message, house) so the caller can render and speak",
])

# ---- 32. edge cases ----
content("devices", "Edge cases we handle", [
    "Empty or gibberish input -> a hint, nothing flips",
    "'everything' means the switchable stuff (light + fan), not the door",
    "Half-written devices.json -> _load() fills in the missing keys",
    "'shut' is in both off-words and close-words - the door branch wins for doors",
    "Two laptops writing at once -> last write wins (acceptable here)",
])

# ---- 33. section: voice ----
section("05", "Voice", "Offline in, offline out, and made to sound human")

# ---- 34. voice goals ----
stat_slide("voice", "The bar we set", [
    ("$0", "no API keys, no quota, ever"),
    ("0", "network calls once models are local"),
    ("~0.5s", "transcribe + reply for a short clip"),
])

# ---- 35. pipeline ----
s = slide()
title(s, "The voice pipeline", "voice")
chain = [("Mic", PAPER_2), ("MediaRecorder\nwebm/opus", CREAM),
         ("POST /api/voice", CREAM), ("ffmpeg\n-> 16k mono PCM", CREAM),
         ("Vosk\nspeech -> text", FOREST), ("intent.py / chatbot", FOREST),
         ("devices.py", CREAM), ("Piper\ntext -> speech", FOREST),
         ("pitch shift\nffmpeg", CREAM), ("browser plays\nthe reply", PAPER_2)]
cols = 5
bw = Inches(2.15); bh = Inches(1.2)
gx = (CONTENT_W - cols * bw) / (cols - 1)
for i, (lbl, fill) in enumerate(chain):
    r, c = divmod(i, cols)
    x = MARGIN + c * (bw + gx)
    y = Inches(2.1) + r * Inches(2.0)
    box(s, x, y, bw, bh, lbl, fill, CREAM if fill == FOREST else INK, 10.5,
        bold=(fill == FOREST))
    if c < cols - 1:
        connect(s, x + bw, y + bh / 2, x + bw + gx, y + bh / 2, w=1.3)
connect(s, MARGIN + 4 * (bw + gx) + bw / 2, Inches(2.1) + bh,
        MARGIN + bw / 2, Inches(4.1), color=INK_SOFT, dash=True)

# ---- 36. Vosk ----
content("voice - in", "Speech to text: Vosk", [
    "A Kaldi-based recognizer with a 40 MB English model - runs on CPU",
    "Streaming-capable, but we feed it the whole clip at once",
    "Output is lower-case, no punctuation - perfect for keyword matching",
    "The small model alone is so-so on free dictation...",
    "...so we constrain it to the words our commands actually use",
])

# ---- 37. accuracy ----
content("voice - accuracy", "Three tricks stacked", [
    "Grammar lock - build a vocabulary from intent.py and let Vosk pick only "
    "from those words. Huge accuracy win for commands.",
    "N-best - ask for the 3 best guesses, take the first that parses as a command",
    "Fuzzy repair - snap near-misses to known words ('fam' -> 'fan') with difflib",
    "Then a free-form pass as a fallback, so the chatbot hook still works",
    "Result: 10/10 commands on clean audio, 12/13 on badly degraded audio",
])

# ---- 38. grammar ----
code_slide("voice.py", "The grammar builds itself from intent.py", '''def _command_grammar():
    words = set(intent.DEVICE_WORDS)
    for s in (intent.TURN_ON_WORDS, intent.TURN_OFF_WORDS,
              intent.OPEN_WORDS, intent.CLOSE_WORDS, intent.EVERYTHING_WORDS):
        words |= set(s)
    words |= {"turn", "switch", "the", "a", "please", "my", "to", ...}
    return json.dumps([" ".join(sorted(words)), "[unk]"])

# two files that can never drift apart - the parser IS the vocabulary''',
    note="SMARTHOME_VOICE_GRAMMAR=0 turns it off. voice_setup.py --big swaps in "
         "a 128 MB model for better free dictation (no grammar).")

# ---- 39. Piper ----
content("voice - out", "Text to speech: Piper", [
    "A neural TTS - the voice is genuinely natural, not the old robotic beep",
    "One 63 MB ONNX voice model, synthesises faster than real time on CPU",
    "We pick a warm US-English voice (lessac, medium quality)",
    "SynthesisConfig tunes pace and the amount of prosody variation",
    "Output is a plain 16-bit WAV the browser plays with no plugin",
])

# ---- 40. humanising ----
code_slide("voice.py", "A touch of pitch for warmth", '''TTS_PITCH = float(os.environ.get("SMARTHOME_TTS_PITCH", "1.06"))

def _pitch_shift(wav_bytes, factor):
    # raise pitch WITHOUT speeding the speech up
    af = f"asetrate={rate}*{factor},aresample={rate},atempo={1/factor:.5f}"
    return ffmpeg(["-i","pipe:0","-af",af,"-f","wav","pipe:1"], wav_bytes)

# 1.00 = Piper as-is,  1.06 = a little warmer / less flat''',
    note="Small shift only. Piper already sounds human; this just lifts it off "
         "the floor. Tunable per deployment.")

# ---- 41. the brain ----
content("voice - brain", "Command first, chat second", [
    "transcribe() -> answer(text, who)",
    "If intent.parse() succeeds -> do it, speak the confirmation ('device')",
    "If not, and SMARTHOME_CHATBOT_URL is set -> ask it, speak the reply ('chat')",
    "Otherwise -> a short 'I can do lights, fan and the door' line ('unknown')",
    "chatbot_reply() is also where a separate assistant's router can drop in",
])

# ---- 42. browser side ----
content("voice - front-end", "What the browser does", [
    "getUserMedia with noise suppression + echo cancel + auto gain",
    "MediaRecorder captures webm/opus",
    "A Web Audio meter watches the level - auto-stops ~1.6s after you finish",
    "10-second hard cap as a backstop; tap the mic again to stop early",
    "A big red 'Listening...' pop-up, then amber 'Thinking...' while it works",
])

# ---- 43. fallback ----
content("voice - resilience", "It degrades, it doesn't break", [
    "GET /api/voice/status tells the page what the server can do",
    "Models + ffmpeg present -> the offline engine",
    "Missing -> the browser's own Web Speech API (needs Chrome + internet)",
    "Neither -> the mic disables itself, typing still works",
    "main.py auto-switches into .venv if it has the voice packages",
])

# ---- 44. setup ----
code_slide("one-time setup", "Getting the offline engine running", '''python -m venv .venv
.venv/bin/pip install -r requirements.txt
.venv/bin/python voice_setup.py          # ~100 MB of models
sudo pacman -S ffmpeg

.venv/bin/python main.py                  # or just: python main.py
#   [main] using .venv/bin/python for offline voice''',
    note="Models land in data/models/ (git-ignored). Everything after is offline.")

# ---- 45. section: front-end ----
section("06", "The front-end", "Vanilla JS, a warm theme, small touches")

# ---- 46. dashboard ----
content("front-end", "The dashboard", [
    "Command panel - mic + a text box, both go to the same place",
    "Three device cards with their own icon, state pill and animation",
    "Quick controls - everything on / off, lock the door",
    "The activity feed",
    "Invites tab for admins, with the pending list and dev-mode toggle",
])

# ---- 47. design language ----
two_col("front-end", "Paper and ink", "The look", [
    "One accent colour - terracotta",
    "Serif body, monospace for labels and log times",
    "Soft shadows, rounded cards",
], "Motion with meaning", [
    "Light flickers on, then glows",
    "Fan blades spin while it runs",
    "The bolt clunks across when the door locks",
    "Room goes dark when the light is off - the page follows",
])

# ---- 48. login touches ----
content("front-end", "Two small touches on the login page", [
    "Arrows scattered on the brand panel that swing to point at your cursor "
    "(atan2 on mousemove, a soft ease so they glide)",
    "A password strength meter on sign-up - weak / fair / good / strong",
    "Score = length + character-class mix, minus points for '1234'-style inputs",
    "Weak passwords are blocked unless developer mode is on",
])

# ---- 49. realtime ----
content("front-end", "Staying in sync", [
    "No websockets - the dashboard polls GET /api/devices every 5 seconds",
    "Renders only what changed, so animations don't re-fire on every poll",
    "The pending list polls every 8 seconds",
    "Toggle a device on one laptop, the other catches up within 5s",
    "Good enough for a demo; a push channel is on the roadmap",
])

# ---- 50. section: quality ----
section("07", "Quality", "Tests, fixes, and honest limitations")

# ---- 51. self-test ----
code_slide("main.py", "The self-test runs on a copy of your data", '''@contextlib.contextmanager
def sandboxed_data():
    tmp = tempfile.mkdtemp(prefix="smarthome-selftest-")
    shutil.copytree("data", os.path.join(tmp, "data"))
    # repoint every module's file paths at the copy...
    yield
    # ...then put them back and delete the copy

# so the test can create users, flip devices, lock accounts - safely''',
    note="python main.py --check runs it and exits. It also runs on every "
         "normal start.")

# ---- 52. what we tested ----
two_col("quality", "What the suite covers", "Green today", [
    "Password scramble round-trips",
    "intent.parse - 5 phrases + 2 rejects",
    "devices.apply / set_device round-trip",
    "request -> approve -> login",
    "voice status + a live STT/TTS round trip",
], "Also checked by hand", [
    "Every API endpoint via curl",
    "Login lockout counter",
    "Voice on clean + degraded audio",
    "Fallback when models are absent",
    "The server starts under both Pythons",
])

# ---- 53. the fix ----
content("quality", "The bug we found on the way", [
    "adminconnection.py - an older, parallel auth system, still in the repo",
    "It used a different users.json layout and single-SHA-256 passwords",
    "Nothing imported it - but running it would rewrite users.json into a "
    "shape login2.py can't read",
    "Moved to archive/ with a note. One landmine defused.",
])

# ---- 54. limitations ----
content("quality", "What it isn't (yet)", [
    "No real hardware - the devices are state in a JSON file",
    "Passwords use a learning-grade scheme, not argon2",
    "Every approved user is a full admin",
    "HTTP only, single machine, no rate limiting",
    "Voice is English-only and command-shaped without the chatbot wired in",
])

# ---- 55. section: roadmap ----
section("08", "Where next", "From prototype to something you'd run at home")

# ---- 56. roadmap ----
table_slide("roadmap", "The next steps", ["Area", "Plan"], [
    ["Hardware", "MQTT / Home Assistant bridge behind devices.py"],
    ["Voice", "A wake word ('hey house') so there's no button"],
    ["Accounts", "Real password hashing, proper roles, HTTPS"],
    ["Realtime", "Server-sent events instead of polling"],
    ["Multi-home", "Namespaced houses, invite links"],
    ["Chatbot", "Wire the assistant router into chatbot_reply()"],
], col_ratios=[1.2, 3.4])

# ---- 57. learned ----
content("close", "What we took away", [
    "The standard library goes further than you'd think",
    "A constrained vocabulary beats a bigger model for a fixed command set",
    "'Degrade, don't break' - every layer has a fallback and the demo survives",
    "Keeping state as readable JSON made debugging almost pleasant",
    "Small touches (a glide on the arrows, a warm pitch) are cheap and land",
])

# ---- 58. try it ----
code_slide("try it", "Five minutes to a talking house", '''git clone <repo> && cd Iron-Titans-dev-strom-public-
python main.py                     # open http://localhost:8000
#   log in:  admin / admin123

# optional - offline voice:
python -m venv .venv && .venv/bin/pip install -r requirements.txt
.venv/bin/python voice_setup.py
python main.py                     # mic now uses Vosk + Piper''',
    note='Then click the mic and say "turn on the light".')

# ---- 59. thanks ----
s = slide(dark=True)
rect(s, 0, 0, Inches(0.18), EMU_H, ACCENT)
text(s, MARGIN, Inches(2.6), CONTENT_W, Inches(1.4),
     "Thank you", size=54, color=CREAM, bold=True, font=HEAD_FONT)
text(s, MARGIN, Inches(3.9), CONTENT_W, Inches(0.6),
     "Questions?  -  and yes, it works offline.", size=18,
     color=RGBColor(0xC7, 0xBE, 0xAA), italic=True)
text(s, MARGIN, Inches(5.6), CONTENT_W, Inches(0.5),
     "Pranjal  -  Samarpreet  -  Ekamjot  -  Amrit", size=13,
     color=RGBColor(0x8A, 0x82, 0x70), font=MONO_FONT)
footer(s, dark=True)

prs.save(OUT)
print(f"wrote {OUT}  ({_num} slides)")
